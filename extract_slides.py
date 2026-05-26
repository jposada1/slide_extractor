"""
Extract slide images and text from student PPTX files.
For full-slide images (Google Slides flattened export), an AI vision model
is used via OpenRouter to extract the chart description and explanation.

Output:
  extracted_image/<pptx_name>/slide001.png, slide002.png, ...
  extracted_slides.csv  (columns: pptx_file, slide_num, img_path, img_type,
                                   explanation, ai_chart, ai_explanation)

  img_type values:
    chart      — image smaller than full slide (an actual chart/plot)
    full_slide — image covers the whole slide (Google Slides flattened export)
    no_image   — slide has no embedded image

Usage:
  pip install python-pptx Pillow openai python-dotenv
  python extract_slides.py
  python extract_slides.py --model google/gemini-2.0-flash-001
"""

import argparse
import base64
import csv
import json
import re
import shutil
from pathlib import Path

from dotenv import load_dotenv
from openai import OpenAI
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import Image
import os

load_dotenv(Path(__file__).parent / ".env")

# ── Paths & config ─────────────────────────────────────────────────────────────

SLIDE_DIR  = Path(__file__).parent
OUTPUT_DIR = SLIDE_DIR / "extracted_image"
CSV_PATH   = SLIDE_DIR / "extracted_slides.csv"

FULL_SLIDE_THRESHOLD = 0.90
MIN_IMAGE_THRESHOLD  = 0.03   # skip images whose area < 3 % of slide area (logos, icons)
DEFAULT_MODEL        = "anthropic/claude-sonnet-4.6"

AI_PROMPT = """You are analyzing a presentation slide image.

First decide: does this slide contain a data visualization (chart, graph, plot, table of data)?
- YES if you can see: bar chart, line chart, scatter plot, pie chart, histogram, heatmap, box plot, data table, etc.
- NO if the slide is only text, a photo, a logo, a title card, a storyboard, or any non-data image.

Return ONLY a JSON object, no markdown:

If the slide does NOT contain a chart/plot:
{"has_chart": false, "plot_bbox": null, "explanation": null}

If the slide DOES contain a chart/plot:
{
  "has_chart": true,
  "plot_bbox": {"x1": <left>, "y1": <top>, "x2": <right>, "y2": <bottom>},
  "explanation": "<all text on the slide: titles, axis labels, captions, bullet points>"
}

plot_bbox must be a tight box around ONLY the chart area (axes + data), expressed as fractions (0.0–1.0) of image width/height.
Do NOT return a box that covers more than 85% of the image."""


# ── OpenRouter client ──────────────────────────────────────────────────────────

def make_client() -> OpenAI:
    key = os.getenv("OPENROUTER_API_KEY", "")
    assert key, "Set OPENROUTER_API_KEY in chart_verifier/.env"
    return OpenAI(api_key=key, base_url="https://openrouter.ai/api/v1")


# ── Utilities ──────────────────────────────────────────────────────────────────

def relative_path(path: Path, anchor: str = "Github") -> str:
    for i, part in enumerate(path.parts):
        if part == anchor:
            return "/".join(path.parts[i:])
    return str(path)


def safe_stem(name: str, max_len: int = 40) -> str:
    return re.sub(r"[^\w]", "_", name)[:max_len]


def save_as_png(blob: bytes, src_ext: str, dest: Path) -> Path:
    dest.write_bytes(blob)
    if src_ext.lower() != "png":
        png = dest.with_suffix(".png")
        with Image.open(dest) as im:
            im.save(png, "PNG")
        dest.unlink()
        return png
    return dest


def encode_image(path: Path) -> str:
    return base64.b64encode(path.read_bytes()).decode()


MIN_TEXT_LENGTH = 30  # characters — shorter strings are treated as titles/labels

def extract_text(slide) -> str:
    shapes = sorted(slide.shapes, key=lambda s: (s.top or 0, s.left or 0))
    lines = []
    for shape in shapes:
        if shape.has_text_frame:
            text = " ".join(shape.text_frame.text.split())
            if len(text) >= MIN_TEXT_LENGTH:
                lines.append(text)
    return " | ".join(lines)


def image_type(shape, slide_w, slide_h) -> str | None:
    w_ratio = shape.width  / slide_w
    h_ratio = shape.height / slide_h
    if w_ratio * h_ratio < MIN_IMAGE_THRESHOLD:
        return None  # too small — likely a logo or icon
    return "full_slide" if (w_ratio >= FULL_SLIDE_THRESHOLD and
                            h_ratio >= FULL_SLIDE_THRESHOLD) else "chart"


# ── AI extraction for full-slide images ───────────────────────────────────────

def crop_chart(img_path: Path, bbox: dict) -> Path:
    """Crop the chart region from the slide image and save as a new PNG."""
    with Image.open(img_path) as im:
        w, h   = im.size
        left   = int(bbox["x1"] * w)
        top    = int(bbox["y1"] * h)
        right  = int(bbox["x2"] * w)
        bottom = int(bbox["y2"] * h)
        cropped = im.crop((left, top, right, bottom))
        chart_path = img_path.with_name(img_path.stem + "_chart.png")
        cropped.save(chart_path, "PNG")
    return chart_path


def ai_extract(img_path: Path, client: OpenAI, model: str) -> dict:
    img_b64 = encode_image(img_path)
    try:
        resp = client.chat.completions.create(
            model=model,
            max_tokens=512,
            messages=[{
                "role": "user",
                "content": [
                    {"type": "text", "text": AI_PROMPT},
                    {"type": "image_url",
                     "image_url": {"url": f"data:image/png;base64,{img_b64}"}},
                ],
            }],
        )
        raw = resp.choices[0].message.content or ""
        obj = json.loads(raw[raw.find("{"):raw.rfind("}") + 1])

        if not obj.get("has_chart", True):
            return {"has_chart": False, "ai_chart_path": None, "ai_explanation": None}

        bbox = obj.get("plot_bbox")
        if bbox:
            bbox_area = (bbox["x2"] - bbox["x1"]) * (bbox["y2"] - bbox["y1"])
            if bbox_area > 0.85:
                bbox = None
        chart_path = crop_chart(img_path, bbox) if bbox else None

        return {
            "has_chart": True,
            "ai_chart_path": relative_path(chart_path) if chart_path else None,
            "ai_explanation": obj.get("explanation"),
        }
    except Exception as e:
        return {"has_chart": True, "ai_chart_path": None, "ai_explanation": f"[error: {e}]"}


# ── Core extraction ────────────────────────────────────────────────────────────

def extract_images(slide, slide_w, slide_h, out_dir: Path, slide_num: int) -> list[dict]:
    results = []
    counter = [1]

    def process(shape):
        itype = image_type(shape, slide_w, slide_h)
        if itype is None:
            return
        ext   = (shape.image.ext or "png").lower().lstrip(".")
        dest  = out_dir / f"slide_{slide_num:03d}_image_{counter[0]:02d}.{ext}"
        saved = save_as_png(shape.image.blob, ext, dest)
        counter[0] += 1
        results.append({
            "img_path":  relative_path(saved),
            "img_type":  itype,
            "_abs_path": saved,
        })

    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            process(shape)
        elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            for child in shape.shapes:
                if child.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    process(child)

    return results


def process_pptx(pptx_path: Path, out_dir: Path,
                 client: OpenAI, model: str) -> list[dict]:
    prs       = Presentation(str(pptx_path))
    slide_out = out_dir / safe_stem(pptx_path.stem)
    if slide_out.exists():
        print(f"  Skipping {pptx_path.name} — folder already exists")
        return []
    slide_out.mkdir(parents=True)

    rows = []
    for num, slide in enumerate(prs.slides, 1):
        images = extract_images(slide, prs.slide_width, prs.slide_height, slide_out, num)
        text   = extract_text(slide)

        if images:
            for img in images:
                if img["img_type"] == "full_slide":
                    print(f"    AI extracting slide {num} ...")
                    ai = ai_extract(img["_abs_path"], client, model)

                    if not ai.get("has_chart", True):
                        print(f"      → no chart, skipping")
                        img["_abs_path"].unlink(missing_ok=True)
                        continue

                    if not ai.get("ai_chart_path"):
                        img["_abs_path"].unlink(missing_ok=True)
                        continue

                    img["_abs_path"].unlink(missing_ok=True)
                    img_path    = ai["ai_chart_path"]
                    explanation = ai.get("ai_explanation") or text
                else:
                    img_path    = img["img_path"]
                    explanation = text

                rows.append({
                    "pptx_file":   pptx_path.name,
                    "slide_num":   num,
                    "img_path":    img_path,
                    "img_type":    "chart",
                    "explanation": explanation,
                })
        # slides with no image are skipped

    return rows


# ── Main ───────────────────────────────────────────────────────────────────────

def main(out_dir: Path, csv_path: Path, model: str) -> None:
    client = make_client()
    out_dir.mkdir(parents=True, exist_ok=True)

    pptx_files = sorted(p for p in SLIDE_DIR.glob("*.pptx") if not p.name.startswith("~$"))
    if not pptx_files:
        print("No .pptx files found in", SLIDE_DIR)
        return

    all_rows = []
    for pptx_path in pptx_files:
        print(f"Processing: {pptx_path.name}")
        rows = process_pptx(pptx_path, out_dir, client, model)
        all_rows.extend(rows)
        n_chart = sum(1 for r in rows if r["img_type"] == "chart")
        n_full  = sum(1 for r in rows if r["img_type"] == "full_slide")
        print(f"  {len(rows)} slides — chart: {n_chart}, full_slide: {n_full}")

    fieldnames = ["pptx_file", "slide_num", "img_path", "img_type", "explanation"]
    with open(csv_path, "w", newline="", encoding="utf-8") as f:
        writer = csv.DictWriter(f, fieldnames=fieldnames)
        writer.writeheader()
        writer.writerows(all_rows)

    print(f"\nDone — {len(all_rows)} rows saved to {csv_path}")


if __name__ == "__main__":
    parser = argparse.ArgumentParser()
    parser.add_argument("--out_dir", type=Path,  default=OUTPUT_DIR)
    parser.add_argument("--csv",     type=Path,  default=CSV_PATH)
    parser.add_argument("--model",   type=str,   default=DEFAULT_MODEL)
    args = parser.parse_args()
    main(args.out_dir, args.csv, args.model)
